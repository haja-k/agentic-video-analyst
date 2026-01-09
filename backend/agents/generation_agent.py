"""
Generation Agent - Handles PDF and PowerPoint generation
"""
from typing import Dict, Any, List, Optional
import logging
from pathlib import Path
from datetime import datetime
import os

from .base_agent import BaseAgent

logger = logging.getLogger(__name__)


class GenerationAgent(BaseAgent):
    """Agent responsible for generating PDF and PPT reports"""
    
    def __init__(self):
        super().__init__("Generation Agent")
        
    async def initialize(self):
        """Initialize generation tools"""
        logger.info("Initializing document generation tools")
        logger.info("✓ Generation tools ready")
        
    async def process(self, input_data: Dict[str, Any]) -> Dict[str, Any]:
        """
        Generate PDF or PowerPoint document
        
        Args:
            input_data: {
                "format": str,  # "pdf" or "pptx"
                "title": str,
                "transcription": Optional[Dict],
                "vision_results": Optional[Dict],
                "summary": Optional[str],
                "output_path": Optional[str]
            }
            
        Returns:
            {
                "status": str,
                "output_path": str,
                "file_size": int
            }
        """
        format_type = input_data.get("format", "pdf")
        output_path = input_data.get("output_path")
        
        if not output_path:
            # Create results directory if it doesn't exist
            results_dir = Path("results")
            results_dir.mkdir(exist_ok=True)
            
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_path = str(results_dir / f"report_{timestamp}.{format_type}")
        else:
            # Ensure the output directory exists when path is provided
            output_dir = Path(output_path).parent
            output_dir.mkdir(parents=True, exist_ok=True)
            
            # Add timestamp and extension if path doesn't have them
            if not output_path.endswith(f".{format_type}"):
                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                output_path = f"{output_path}_{timestamp}.{format_type}"
        
        logger.info(f"Generating {format_type.upper()} document: {output_path}")
        
        # Extract the actual content dict
        content = input_data.get("content", input_data)
        
        try:
            if format_type == "pdf":
                return await self.generate_pdf(content, output_path)
            elif format_type == "pptx":
                return await self.generate_pptx(content, output_path)
            else:
                return {"error": f"Unsupported format: {format_type}"}
        except Exception as e:
            logger.error(f"Generation failed: {e}")
            return {"error": str(e)}
    
    async def generate_pdf(self, content: Dict, output_path: str) -> Dict[str, Any]:
        """Generate PDF document using ReportLab"""
        from reportlab.lib.pagesizes import letter, A4
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import inch
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
        from reportlab.platypus import Table, TableStyle
        from reportlab.lib import colors
        
        # Debug: Log what content we received
        logger.info(f"PDF content keys: {content.keys()}")
        logger.info(f"Transcription present: {bool(content.get('transcription'))}")
        logger.info(f"Vision results present: {bool(content.get('vision_results'))}")
        
        doc = SimpleDocTemplate(output_path, pagesize=letter)
        styles = getSampleStyleSheet()
        story = []
        
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=24,
            textColor=colors.HexColor('#1a1a1a'),
            spaceAfter=30
        )
        
        title = content.get("title", "Video Analysis Report")
        story.append(Paragraph(title, title_style))
        story.append(Spacer(1, 0.2*inch))
        
        story.append(Paragraph(f"Generated: {datetime.now().strftime('%B %d, %Y at %H:%M')}", styles['Normal']))
        story.append(Spacer(1, 0.5*inch))
        
        if content.get("summary"):
            story.append(Paragraph("Executive Summary", styles['Heading2']))
            story.append(Spacer(1, 0.1*inch))
            story.append(Paragraph(content["summary"], styles['Normal']))
            story.append(Spacer(1, 0.3*inch))
        
        if content.get("transcription"):
            trans = content["transcription"]
            story.append(Paragraph("Transcription", styles['Heading2']))
            story.append(Spacer(1, 0.1*inch))
            
            if isinstance(trans, dict) and "transcription" in trans:
                story.append(Paragraph(trans["transcription"], styles['Normal']))
                
                if trans.get("segments"):
                    story.append(Spacer(1, 0.2*inch))
                    story.append(Paragraph("Key Segments", styles['Heading3']))
                    story.append(Spacer(1, 0.1*inch))
                    
                    for seg in trans["segments"][:10]:
                        timestamp = f"[{seg['start']:.1f}s - {seg['end']:.1f}s]"
                        story.append(Paragraph(f"<b>{timestamp}</b> {seg['text']}", styles['Normal']))
                        story.append(Spacer(1, 0.05*inch))
            
            story.append(Spacer(1, 0.3*inch))
        
        if content.get("vision_results"):
            vision = content["vision_results"]
            story.append(Paragraph("Visual Analysis", styles['Heading2']))
            story.append(Spacer(1, 0.1*inch))
            
            if isinstance(vision, dict) and vision.get("results"):
                story.append(Paragraph(f"Analyzed {vision.get('frames_analyzed', 0)} frames", styles['Normal']))
                story.append(Spacer(1, 0.2*inch))
                
                for i, frame in enumerate(vision["results"][:5], 1):
                    story.append(Paragraph(f"Frame {i} (at {frame.get('timestamp', 0):.1f}s)", styles['Heading3']))
                    
                    if frame.get("caption"):
                        story.append(Paragraph(f"Scene: {frame['caption']}", styles['Normal']))
                    
                    if frame.get("objects"):
                        objects_text = ", ".join([obj['class'] for obj in frame['objects'][:5]])
                        story.append(Paragraph(f"Objects detected: {objects_text}", styles['Normal']))
                    
                    story.append(Spacer(1, 0.15*inch))
        
        doc.build(story)
        
        file_size = os.path.getsize(output_path)
        logger.info(f"PDF generated: {output_path} ({file_size} bytes)")
        
        return {
            "status": "success",
            "output_path": output_path,
            "file_size": file_size
        }
    
    async def generate_pptx(self, content: Dict, output_path: str) -> Dict[str, Any]:
        """Generate PowerPoint presentation using python-pptx"""
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = content.get("title", "Video Analysis Report")
        subtitle.text = f"Generated on {datetime.now().strftime('%B %d, %Y')}"
        
        # Format title slide
        title.text_frame.paragraphs[0].font.name = 'Calibri'
        title.text_frame.paragraphs[0].font.size = Pt(20)
        subtitle.text_frame.paragraphs[0].font.name = 'Calibri'
        subtitle.text_frame.paragraphs[0].font.size = Pt(17)
        
        if content.get("summary"):
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            title = slide.shapes.title
            body = slide.placeholders[1]
            
            title.text = "Executive Summary"
            tf = body.text_frame
            tf.text = content["summary"]
            
            # Format text
            title.text_frame.paragraphs[0].font.name = 'Calibri'
            title.text_frame.paragraphs[0].font.size = Pt(20)
            for paragraph in tf.paragraphs:
                paragraph.font.name = 'Calibri'
                paragraph.font.size = Pt(15)
        
        if content.get("transcription"):
            trans = content["transcription"]
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            title = slide.shapes.title
            body = slide.placeholders[1]
            
            title.text = "Transcription"
            tf = body.text_frame
            
            # Format title
            title.text_frame.paragraphs[0].font.name = 'Calibri'
            title.text_frame.paragraphs[0].font.size = Pt(28)
            
            if isinstance(trans, dict):
                if trans.get("transcription"):
                    text = trans["transcription"]
                    if len(text) > 500:
                        text = text[:500] + "..."
                    tf.text = text
                    # Format text
                    for paragraph in tf.paragraphs:
                        paragraph.font.name = 'Calibri'
                        paragraph.font.size = Pt(15)
                
                if trans.get("segments"):
                    slide = prs.slides.add_slide(bullet_slide_layout)
                    title = slide.shapes.title
                    body = slide.placeholders[1]
                    title.text = "Key Segments"
                    tf = body.text_frame
                    
                    # Format title
                    title.text_frame.paragraphs[0].font.name = 'Calibri'
                    title.text_frame.paragraphs[0].font.size = Pt(20)
                    
                    for i, seg in enumerate(trans["segments"][:8]):
                        if i == 0:
                            tf.text = f"{seg['start']:.1f}s: {seg['text'][:80]}"
                            tf.paragraphs[0].font.name = 'Calibri'
                            tf.paragraphs[0].font.size = Pt(15)
                        else:
                            p = tf.add_paragraph()
                            p.text = f"{seg['start']:.1f}s: {seg['text'][:80]}"
                            p.level = 0
                            p.font.name = 'Calibri'
                            p.font.size = Pt(15)
        
        if content.get("vision_results"):
            vision = content["vision_results"]
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            title = slide.shapes.title
            body = slide.placeholders[1]
            
            title.text = "Visual Analysis"
            
            # Format title
            title.text_frame.paragraphs[0].font.name = 'Calibri'
            title.text_frame.paragraphs[0].font.size = Pt(28)
            
            if isinstance(vision, dict) and vision.get("results"):
                tf = body.text_frame
                tf.text = f"Analyzed {vision.get('frames_analyzed', 0)} frames from the video"
                tf.paragraphs[0].font.name = 'Calibri'
                tf.paragraphs[0].font.size = Pt(15)
                
                for i, frame in enumerate(vision["results"][:4], 1):
                    p = tf.add_paragraph()
                    timestamp = frame.get('timestamp', 0)
                    caption = frame.get('caption', 'No description')
                    p.text = f"Frame at {timestamp:.1f}s: {caption}"
                    p.level = 1
                    p.font.name = 'Calibri'
                    p.font.size = Pt(15)
                    
                    if frame.get("objects"):
                        obj_names = [obj['class'] for obj in frame['objects'][:3]]
                        p2 = tf.add_paragraph()
                        p2.text = f"Objects: {', '.join(obj_names)}"
                        p2.level = 2
                        p2.font.name = 'Calibri'
                        p2.font.size = Pt(15)
        
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        left = Inches(2)
        top = Inches(3)
        width = Inches(6)
        height = Inches(1)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = "Thank You"
        p = tf.paragraphs[0]
        p.font.name = 'Calibri'
        p.font.size = Pt(30)
        p.alignment = PP_ALIGN.CENTER
        
        prs.save(output_path)
        
        file_size = os.path.getsize(output_path)
        logger.info(f"PowerPoint generated: {output_path} ({file_size} bytes)")
        
        return {
            "status": "success",
            "output_path": output_path,
            "file_size": file_size
        }
    
    def format_content_for_pdf(self, content: Dict) -> List[Dict]:
        """Format content structure for PDF layout"""
        # TODO: Structure content into sections, paragraphs, lists
        pass
    
    def format_content_for_pptx(self, content: Dict) -> List[Dict]:
        """Format content structure for PowerPoint slides"""
        # TODO: Structure content into slides with titles, bullets, images
        pass
